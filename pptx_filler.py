# -*- coding: utf-8 -*-
"""
Logica di compilazione del PowerPoint (indipendente da Streamlit, così è
testabile da sola).

Funzioni principali:
  - load_mapping(xlsx, regenerate=False) -> dict {segnaposto: testo}
  - extract_placeholders(pptx)           -> set di segnaposto presenti nel file
  - fill_pptx(pptx, mapping)             -> (BytesIO compilato, stats)

Gestisce in modo robusto:
  - segnaposto {{chiave}} spezzati su più "run" dentro lo stesso paragrafo
  - testo dentro tabelle
  - forme raggruppate (group shapes), in modo ricorsivo
"""
import io
import math
import re
import unicodedata
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# backend headless per generare il radar anche senza display
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as _plt
    _HAS_MPL = True
except Exception:                       # matplotlib assente -> radar saltato
    _HAS_MPL = False

PLACEHOLDER_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")

# le 4 aree dell'assessment, nell'ordine usato dal radar
AREA_ORDER = ("strategia", "cultura", "processi", "mercato")

# codici domanda esclusi dal calcolo degli score (peso nullo nella mappatura
# originale). Dal foglio "Metodologia calcolo": "PQ19 è presente nella survey
# e nella mappatura ma non ha peso PMI/GRANDI" -> va esclusa.
_PESO_ZERO_COD = {"PQ19"}


# ----------------------------------------------------------------- Excel -> mapping
def _records(ws):
    """Legge un foglio come lista di dict, header in prima riga (lower/strip)."""
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return []
    headers = [(str(h).strip().lower() if h is not None else "") for h in rows[0]]
    out = []
    for row in rows[1:]:
        if all(c is None or str(c).strip() == "" for c in row):
            continue
        out.append({h: row[i] for i, h in enumerate(headers) if h})
    return out


def _read_scala(wb):
    """{livello: (nome, descrizione)} dal foglio 'Scala', se presente."""
    scala = {}
    if "Scala" in wb.sheetnames:
        for rec in _records(wb["Scala"]):
            lv = rec.get("livello")
            if lv is None:
                continue
            nome = str(rec.get("nome", "")).strip()
            desc = "" if rec.get("descrizione") is None else str(rec["descrizione"]).strip()
            scala[int(lv)] = (nome, desc)
    return scala


def _derive_panoramica(mapping, scala):
    """Inietta i segnaposto panoramica.* : il livello macro aggregato
    (media dei livelli delle 4 aree). scala: {livello: (nome, descrizione)}.

    Le aree si riconoscono perché l'id NON contiene '_' (i sottogruppi sì),
    quindi la funzione è indipendente dalla fonte del mapping.
    """
    livelli = []
    for key, val in mapping.items():
        if not key.endswith(".risultato"):
            continue
        aid = key[: -len(".risultato")]
        if "_" in aid:                       # salta i sottogruppi
            continue
        m = re.match(r"\s*(\d+)", str(val))
        if m:
            livelli.append(int(m.group(1)))
    if not livelli or not scala:
        return
    avg = sum(livelli) / len(livelli)
    macro = max(1, min(4, int(avg + 0.5)))   # arrotondamento half-up, clamp 1..4
    nome = scala.get(macro, ("", ""))[0] or str(macro)
    desc = scala.get(macro, ("", ""))[1] or ""
    mapping["panoramica.risultato"] = f"{macro} - {nome}"
    mapping["panoramica.livello_medio"] = f"{avg:.2f}".replace(".", ",") + " / 4"
    mapping["panoramica.descrizione"] = desc


def _area_scores(mapping):
    """[(nome, livello)] delle 4 aree, ricavati dai <area>.risultato nel mapping."""
    scores = []
    for aid in AREA_ORDER:
        m = re.match(r"\s*(\d+)", str(mapping.get(f"{aid}.risultato", "")))
        scores.append((aid.upper(), int(m.group(1)) if m else 0))
    return scores


# ----------------------------------------------------- survey Qualtrics (reale)
# Soglie ufficiali score -> livello (dal foglio "Metodologia calcolo"):
#   1 (Deriva)   1,00-1,50
#   2 (Rotta)    1,51-2,50
#   3 (Orbita)   2,51-3,50
#   4 (Gravità)  > 3,50
def _score_to_level(score):
    if score <= 1.50:
        return 1
    if score <= 2.50:
        return 2
    if score <= 3.50:
        return 3
    return 4


def _norm_label(s):
    """Normalizza per match case/accenti/punteggiatura insensibile."""
    s = unicodedata.normalize("NFKD", str(s))
    s = "".join(c for c in s if not unicodedata.combining(c))
    s = s.lower()
    return re.sub(r"[^a-z0-9]", "", s)


def _libreria_lookup(wb_map):
    """{norm(nome): (id, tipo)} dalla Libreria, per ricondurre le colonne
    'Score ...' della survey agli id usati dai segnaposto."""
    out = {}
    if "Libreria" not in wb_map.sheetnames:
        return out
    for rec in _records(wb_map["Libreria"]):
        _id = str(rec.get("id", "")).strip()
        nome = str(rec.get("nome", "")).strip()
        tipo = str(rec.get("tipo", "")).strip().lower()
        if nome and _id:
            out[_norm_label(nome)] = (_id, tipo)
    return out


def _has_score_columns(wb):
    """True se un foglio ha colonne che iniziano con 'Score ' (formato Qualtrics)."""
    for name in wb.sheetnames:
        hdr = next(wb[name].iter_rows(values_only=True), None)
        if hdr and any(str(c).strip().lower().startswith("score ") for c in hdr if c):
            return True
    return False


def mapping_from_qualtrics(input_wb, mapping_xlsx):
    """Mapping dalla survey Qualtrics reale.

    Colonne attese: 'Score <area>' e 'Score <area> - <sottogruppo>'.
    Aggrega per MEDIA degli score su tutti i rispondenti, deriva il livello
    con le soglie ufficiali (_score_to_level), poi costruisce i segnaposto
    <id>.risultato (+ .commento per le aree / .attivita per i sottogruppi) e
    la panoramica aggregata sulla media degli score delle 4 aree.
    """
    wb_map = load_workbook(mapping_xlsx, data_only=True)
    scala = _read_scala(wb_map)                       # {liv: (nome, desc)}

    testo_by_id_lv = {}                               # (id, liv) -> testo
    if "Libreria" in wb_map.sheetnames:
        for rec in _records(wb_map["Libreria"]):
            _id = str(rec.get("id", "")).strip()
            lv = rec.get("livello")
            if not _id or lv is None:
                continue
            testo_by_id_lv[(_id, int(lv))] = (
                "" if rec.get("testo") is None else str(rec["testo"])
            )
    lookup = _libreria_lookup(wb_map)                 # norm(nome) -> (id, tipo)

    # --- trova il foglio con le colonne 'Score ...' ---
    ws = None
    for name in input_wb.sheetnames:
        cand = input_wb[name]
        hdr = next(cand.iter_rows(values_only=True), ())
        hdr = [str(c).strip() if c is not None else "" for c in hdr]
        if any(h.lower().startswith("score ") for h in hdr):
            ws = cand
            break
    if ws is None:
        raise ValueError("Nessun foglio con colonne 'Score ...' nella survey.")

    rows = list(ws.iter_rows(values_only=True))
    headers = [str(c).strip() if c is not None else "" for c in rows[0]]

    # --- mappa ogni colonna Score -> (id, tipo) tramite il nome su Libreria ---
    score_cols = {}                                   # col_idx -> (id, tipo)
    for ci, h in enumerate(headers):
        if not h.lower().startswith("score "):
            continue
        label = h[len("score "):].strip()
        sub = label.split(" - ", 1)[1].strip() if " - " in label else label
        hit = lookup.get(_norm_label(sub))
        if hit:
            score_cols[ci] = hit
    tipo_by_sid = {sid: tipo for sid, tipo in score_cols.values()}

    # --- media degli score per id su tutti i rispondenti (skip righe metadati) ---
    sums = {}                                         # id -> [score...]
    for row in rows[1:]:
        for ci, (sid, _tipo) in score_cols.items():
            v = row[ci] if ci < len(row) else None
            if isinstance(v, (int, float)) and not isinstance(v, bool):
                sums.setdefault(sid, []).append(float(v))

    mapping = {}
    area_scores = {}                                  # id area -> score medio
    for sid, vals in sums.items():
        mean = sum(vals) / len(vals)
        lv = _score_to_level(mean)
        nome_liv = scala.get(lv, (str(lv), ""))[0]
        mapping[f"{sid}.risultato"] = f"{lv} - {nome_liv}"
        campo = "commento" if tipo_by_sid.get(sid) == "area" else "attivita"
        mapping[f"{sid}.{campo}"] = testo_by_id_lv.get((sid, lv), "")
        if tipo_by_sid.get(sid) == "area":
            area_scores[sid] = mean

    # --- panoramica aggregata: media degli score delle 4 aree -> livello macro ---
    if area_scores and scala:
        overall = sum(area_scores.values()) / len(area_scores)
        macro = _score_to_level(overall)
        nome = scala.get(macro, ("", ""))[0] or str(macro)
        desc = scala.get(macro, ("", ""))[1] or ""
        mapping["panoramica.risultato"] = f"{macro} - {nome}"
        mapping["panoramica.livello_medio"] = (
            f"{overall:.2f}".replace(".", ",") + " / 4"
        )
        mapping["panoramica.descrizione"] = desc
    return mapping


# ----------------------------------------------------- survey Qualtrics RAW
def _read_mappatura(wb_map):
    """{cod: (categoria, sottogruppo, peso_pmi, peso_grandi)} dal foglio dei
    pesi ('Mappatura' oppure 'Mappatura utilizzata'). Vuoto se assente."""
    name = None
    for cand in ("Mappatura", "Mappatura utilizzata"):
        if cand in wb_map.sheetnames:
            name = cand
            break
    if name is None:
        return {}
    out = {}

    def _num(rec, key):
        v = rec.get(key)
        return v if isinstance(v, (int, float)) and not isinstance(v, bool) else None

    for rec in _records(wb_map[name]):
        cod = str(rec.get("cod.", "") or rec.get("cod", "")).strip()
        if not cod:
            continue
        out[cod] = (
            str(rec.get("categoria", "")).strip(),
            str(rec.get("sottogruppo", "")).strip(),
            _num(rec, "peso pmi"),
            _num(rec, "peso grandi"),
        )
    return out


def _a3_to_tipo(a3):
    """PMI / GRANDI dalla fascia di dipendenti (colonna A3).
    Metodologia: PMI se classe sotto i 250 dipendenti, GRANDI altrimenti."""
    s = str(a3).lower() if a3 is not None else ""
    if any(k in s for k in ("250", "500", "1000", "più", "oltre")):
        return "GRANDI"
    return "PMI"


def _has_raw_questions(wb):
    """True se la survey ha colonne domanda CQ/PQ/SQ/MQ e NESSUNA colonna
    'Score ' precalcolata (formato raw da ricalcolare)."""
    if _has_score_columns(wb):
        return False
    for name in wb.sheetnames:
        hdr = next(wb[name].iter_rows(values_only=True), None)
        if hdr and any(re.match(r"^(CQ|PQ|SQ|MQ)\d", str(c).strip())
                       for c in hdr if c):
            return True
    return False


def mapping_from_raw(raw_wb, mapping_xlsx):
    """Mapping dalla survey Qualtrics RAW (solo risposte alle domande, senza
    colonne Score/Livello già calcolate).

    Ricalcola gli score pesati (score = Σ risposta×peso / Σ peso) usando la
    mappatura domanda->area/sottogruppo + Peso PMI/GRANDI; il Tipo impresa
    (PMI/GRANDI) è derivato dalla fascia di dipendenti (A3). PQ19 esclusa
    (vedi _PESO_ZERO_COD). Più rispondenti aggregati per media.
    """
    wb_map = load_workbook(mapping_xlsx, data_only=True)
    scala = _read_scala(wb_map)

    testo_by_id_lv = {}                               # (id, liv) -> testo
    if "Libreria" in wb_map.sheetnames:
        for rec in _records(wb_map["Libreria"]):
            _id = str(rec.get("id", "")).strip()
            lv = rec.get("livello")
            if not _id or lv is None:
                continue
            testo_by_id_lv[(_id, int(lv))] = (
                "" if rec.get("testo") is None else str(rec["testo"])
            )
    lookup = _libreria_lookup(wb_map)                 # norm(nome) -> (id, tipo)
    mapp = _read_mappatura(wb_map)                    # cod -> (cat, sub, pmi, grand)
    if not mapp:
        raise ValueError(
            "Mappatura pesi mancante: nel file DEIA serve un foglio 'Mappatura' "
            "(Cod.|Categoria|Sottogruppo|Peso PMI|Peso GRANDI) per ricalcolare "
            "la survey raw."
        )

    # --- trova il foglio con le colonne domanda ---
    ws = None
    for name in raw_wb.sheetnames:
        hdr = next(raw_wb[name].iter_rows(values_only=True), ())
        hdr_s = [str(c).strip() if c is not None else "" for c in hdr]
        if any(re.match(r"^(CQ|PQ|SQ|MQ)\d", h) for h in hdr_s):
            ws = raw_wb[name]
            break
    if ws is None:
        raise ValueError("Nessun foglio con domande CQ/PQ/SQ/MQ nella survey raw.")
    rows = list(ws.iter_rows(values_only=True))
    headers = [str(c).strip() if c is not None else "" for c in rows[0]]
    idx = {}
    for i, h in enumerate(headers):
        idx.setdefault(h, i)
    a3i = idx.get("A3")

    # cod domanda -> indice colonna
    qcol = {}
    for cod in mapp:
        ci = idx.get(cod)
        if ci is None:                                # match tollerante sugli spazi
            for h, i in idx.items():
                if h.replace(" ", "") == cod.replace(" ", ""):
                    ci = i
                    break
        if ci is not None:
            qcol[cod] = ci

    def _num(v):
        if isinstance(v, bool):
            return None
        if isinstance(v, (int, float)):
            return float(v)
        try:
            return float(str(v).strip())
        except (ValueError, TypeError):
            return None

    # rispondenti = righe con almeno una risposta numerica (salta header/testo)
    resp_rows = []
    for r in rows[1:]:
        if any(_num(r[ci]) is not None for ci in qcol.values() if ci < len(r)):
            resp_rows.append(r)

    area_acc, sub_acc = {}, {}                        # id -> [score per rispondente]
    for r in resp_rows:
        tipo = _a3_to_tipo(r[a3i]) if a3i is not None and a3i < len(r) else "PMI"
        by_cat, by_sub = {}, {}
        for cod, ci in qcol.items():
            if cod in _PESO_ZERO_COD:
                continue
            cat, sub, pmi, grand = mapp[cod]
            peso = pmi if tipo == "PMI" else grand
            if peso is None:
                continue
            a = _num(r[ci]) if ci < len(r) else None
            if a is None:
                continue
            nc, dc = by_cat.get(cat, (0.0, 0.0))
            by_cat[cat] = (nc + a * peso, dc + peso)
            key = (cat, sub)
            ns, ds = by_sub.get(key, (0.0, 0.0))
            by_sub[key] = (ns + a * peso, ds + peso)
        for cat, (n_, d_) in by_cat.items():
            aid = lookup.get(_norm_label(cat), (None, None))[0]
            if aid and d_:
                area_acc.setdefault(aid, []).append(n_ / d_)
        for (cat, sub), (n_, d_) in by_sub.items():
            sid = lookup.get(_norm_label(sub), (None, None))[0]
            if sid and d_:
                sub_acc.setdefault(sid, []).append(n_ / d_)

    mapping = {}
    area_scores = {}
    for aid, vals in area_acc.items():
        mean = sum(vals) / len(vals)
        area_scores[aid] = mean
        lv = _score_to_level(mean)
        nome_liv = scala.get(lv, (str(lv), ""))[0]
        mapping[f"{aid}.risultato"] = f"{lv} - {nome_liv}"
        mapping[f"{aid}.commento"] = testo_by_id_lv.get((aid, lv), "")
    for sid, vals in sub_acc.items():
        mean = sum(vals) / len(vals)
        lv = _score_to_level(mean)
        nome_liv = scala.get(lv, (str(lv), ""))[0]
        mapping[f"{sid}.risultato"] = f"{lv} - {nome_liv}"
        mapping[f"{sid}.attivita"] = testo_by_id_lv.get((sid, lv), "")

    if area_scores and scala:
        overall = sum(area_scores.values()) / len(area_scores)
        macro = _score_to_level(overall)
        nome = scala.get(macro, ("", ""))[0] or str(macro)
        desc = scala.get(macro, ("", ""))[1] or ""
        mapping["panoramica.risultato"] = f"{macro} - {nome}"
        mapping["panoramica.livello_medio"] = f"{overall:.2f}".replace(".", ",") + " / 4"
        mapping["panoramica.descrizione"] = desc
    return mapping


def _mapping_from_sheet(ws):
    mapping = {}
    for rec in _records(ws):
        key = rec.get("segnaposto")
        if key is None:
            continue
        key = str(key).strip()
        if key:
            mapping[key] = "" if rec.get("testo") is None else str(rec["testo"])
    return mapping


def _mapping_from_levels(wb):
    """Ricostruisce il mapping da Assessment + Libreria + Scala (dinamico)."""
    scala = _read_scala(wb)  # {livello: (nome, descrizione)}

    libreria = {}  # (id, livello) -> testo ; e tipo per id
    tipo_by_id = {}
    for rec in _records(wb["Libreria"]):
        _id = str(rec.get("id", "")).strip()
        lv = rec.get("livello")
        if not _id or lv is None:
            continue
        libreria[(_id, int(lv))] = "" if rec.get("testo") is None else str(rec["testo"])
        tipo_by_id[_id] = str(rec.get("tipo", "")).strip().lower()

    mapping = {}
    for rec in _records(wb["Assessment"]):
        _id = str(rec.get("id", "")).strip()
        lv = rec.get("livello")
        if not _id or lv is None:
            continue
        lv = int(lv)
        nome_liv = scala.get(lv, (str(lv), ""))[0]
        mapping[f"{_id}.risultato"] = f"{lv} - {nome_liv}"
        tipo = tipo_by_id.get(_id, "")
        campo = "commento" if tipo == "area" else "attivita"
        mapping[f"{_id}.{campo}"] = libreria.get((_id, lv), "")
    _derive_panoramica(mapping, scala)
    return mapping


def load_mapping(xlsx, regenerate=False):
    """xlsx: path o file-like. regenerate=True ricostruisce da Assessment+Libreria."""
    wb = load_workbook(xlsx, data_only=True)
    if not regenerate and "Mapping" in wb.sheetnames:
        m = _mapping_from_sheet(wb["Mapping"])
        if m:
            _derive_panoramica(m, _read_scala(wb))
            return m
    if "Assessment" in wb.sheetnames and "Libreria" in wb.sheetnames:
        return _mapping_from_levels(wb)            # panoramica già inclusa
    if "Mapping" in wb.sheetnames:
        m = _mapping_from_sheet(wb["Mapping"])
        _derive_panoramica(m, _read_scala(wb))
        return m
    raise ValueError(
        "Excel non valido: serve un foglio 'Mapping' (segnaposto|testo) "
        "oppure i fogli 'Assessment' + 'Libreria'."
    )


def mapping_from_survey(input_xlsx, mapping_xlsx):
    """
    Costruisce il mapping dai risultati della survey.

    Riconosce automaticamente tre formati di input:
      - Qualtrics RAW: solo risposte alle domande CQ/PQ/SQ/MQ (nessuno Score
        precalcolato). Score ricalcolati con pesi PMI/GRANDI (vedi
        mapping_from_raw). Esempio: 'Galactica Prova_...xlsx'.
      - Qualtrics con Score: colonne 'Score <area>' / 'Score <area> - <sottogruppo>'
        con più rispondenti, aggregati per media (vedi mapping_from_qualtrics).
        Esempio: 'Galactica_Prova_calcolo score.xlsx'.

    mapping_xlsx: deia_mapping_GM.xlsx con i fogli Libreria, Scala e (per la
    survey raw) Mappatura.
    """
    wb_in = load_workbook(input_xlsx, data_only=True)

    # --- formato Qualtrics RAW (solo risposte, score da ricalcolare) ---------
    if _has_raw_questions(wb_in):
        return mapping_from_raw(wb_in, mapping_xlsx)

    # --- formato Qualtrics con Score già calcolati ---------------------------
    if _has_score_columns(wb_in):
        return mapping_from_qualtrics(wb_in, mapping_xlsx)

    raise ValueError(
        "Formato survey non riconosciuto. Serve la survey Qualtrics raw "
        "(colonne CQ/PQ/SQ/MQ) oppure con Score gia' calcolati "
        "(colonne 'Score <area>')."
    )


# ----------------------------------------------------------------- PPTX traversal
def _iter_shapes(shapes):
    """Itera ricorsivamente tutte le forme, entrando nei gruppi."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _iter_text_frames(prs):
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame


def _apply_to_text(text, mapping, used, unresolved):
    def repl(m):
        key = m.group(1).strip()
        if key in mapping:
            used.add(key)
            return str(mapping[key])
        unresolved.add(key)
        return m.group(0)
    return PLACEHOLDER_RE.sub(repl, text)


def _replace_in_text_frame(tf, mapping, used, unresolved):
    for para in tf.paragraphs:
        runs = para.runs
        full = "".join(r.text for r in runs)
        if "{{" not in full:
            continue
        new = _apply_to_text(full, mapping, used, unresolved)
        if new == full:
            continue
        if runs:
            runs[0].text = new          # tutto il testo nel primo run
            for r in runs[1:]:
                r.text = ""             # svuota i run successivi
        else:
            para.text = new


# --------------------------------------------------------- radar dinamico
ORANGE_HEX = "#E86A17"
DARK_HEX = "#1A1A1A"
RADAR_RE = re.compile(r"\{\{\s*radar\.aree\s*\}\}")


def _radar_image(scores):
    """scores: [(label, valore 0..4), ...]. Ritorna BytesIO di un PNG radar."""
    labels = [s[0] for s in scores]
    vals = [min(s[1], 4) for s in scores]
    n = len(labels)
    angles = [i / n * 2 * math.pi for i in range(n)]
    angles_c = angles + angles[:1]
    values_c = vals + vals[:1]

    fig = _plt.figure(figsize=(6, 6), dpi=150)
    ax = fig.add_subplot(111, polar=True)
    ax.set_theta_offset(math.pi / 2)
    ax.set_theta_direction(-1)
    ax.set_ylim(0, 4)
    ax.set_yticks([1, 2, 3, 4])
    ax.set_yticklabels(["1", "2", "3", "4"], fontsize=9, color="#999999")
    ax.set_xticks(angles)
    ax.set_xticklabels(labels, fontsize=12, fontweight="bold", color=DARK_HEX)
    ax.tick_params(axis="x", pad=16)
    ax.plot(angles_c, values_c, color=ORANGE_HEX, linewidth=2.2)
    ax.fill(angles_c, values_c, color=ORANGE_HEX, alpha=0.22)
    ax.scatter(angles, vals, color=ORANGE_HEX, s=35, zorder=5)
    for ang, val in zip(angles, vals):
        ax.text(ang, val + 0.28, str(val), color=ORANGE_HEX,
                fontsize=10, fontweight="bold", ha="center", va="center")
    ax.spines["polar"].set_color("#CCCCCC")
    ax.grid(color="#DDDDDD")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.2)
    _plt.close(fig)
    buf.seek(0)
    return buf


def _install_dynamic_images(prs, mapping):
    """Sostituisce i marcatori {{radar.aree}} con un'immagine radar generata
    dai punteggi delle aree presenti nel mapping."""
    if not _HAS_MPL:
        return
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            full = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
            if not RADAR_RE.search(full):
                continue
            img = _radar_image(_area_scores(mapping))
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            size = min(w, h)                          # radar quadrato, centrato
            left = l + (w - size) // 2
            top = t + (h - size) // 2
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(img, left, top, size, size)


def extract_placeholders(pptx):
    prs = Presentation(pptx)
    found = set()
    for tf in _iter_text_frames(prs):
        for para in tf.paragraphs:
            full = "".join(r.text for r in para.runs)
            for m in PLACEHOLDER_RE.finditer(full):
                key = m.group(1).strip()
                if key.startswith("radar"):          # marcatore immagine, non testo
                    continue
                found.add(key)
    return found


def fill_pptx(pptx, mapping):
    """Compila il pptx. Ritorna (BytesIO, stats)."""
    prs = Presentation(pptx)
    _install_dynamic_images(prs, mapping)            # prima le immagini dinamiche
    used, unresolved = set(), set()
    for tf in _iter_text_frames(prs):
        _replace_in_text_frame(tf, mapping, used, unresolved)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    stats = {
        "sostituiti": sorted(used),
        "non_risolti": sorted(unresolved),
        "n_sostituiti": len(used),
        "n_non_risolti": len(unresolved),
        "inutilizzati": sorted(set(mapping) - used),
    }
    return buf, stats
