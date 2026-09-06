# -*- coding: utf-8 -*-
"""Verifica end-to-end della pipeline app (senza Streamlit)."""
import io
import json
import sys
from pathlib import Path

# console Windows (cp1252): forza UTF-8 per stampare emoji/accenti
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt
from mod1.pptx_filler import load_mapping, extract_placeholders, fill_pptx

BASE = Path(__file__).resolve().parent
ok = True


def check(cond, msg):
    global ok
    print(("  OK  " if cond else " FAIL ") + msg)
    if not cond:
        ok = False


def _save(prs):
    b = io.BytesIO()
    prs.save(b)
    b.seek(0)
    return b


def all_text(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        out.append(c.text)
    return "\n".join(out)


print("== 1. Mapping da Excel (default) ==")
mapping = load_mapping(BASE / "deia_mapping_GM.xlsx")
check(len(mapping) == 35, f"35 voci di mapping (trovate {len(mapping)})")
check(mapping.get("strategia.risultato") == "2 - ROTTA",
      "strategia.risultato == '2 - ROTTA'")
check(mapping.get("cultura.risultato") == "1 - DERIVA",
      "cultura.risultato == '1 - DERIVA'")
check("inizia a definire" in mapping.get("strategia.commento", ""),
      "strategia.commento = testo livello 2 (dal documento)")
check(mapping.get("panoramica.risultato") == "1 - DERIVA",
      "panoramica.risultato == '1 - DERIVA' (livello macro)")
check(mapping.get("panoramica.livello_medio") == "1,25 / 4",
      "panoramica.livello_medio == '1,25 / 4'")

print("== 2. Mapping rigenerato dai livelli (Assessment+Libreria) ==")
mapping2 = load_mapping(BASE / "deia_mapping_GM.xlsx", regenerate=True)
check(mapping2 == mapping, "mapping rigenerato identico a quello pronto")

print("== 3. Segnaposto nel template ==")
ph = extract_placeholders(BASE / "template_esempio.pptx")
check(len(ph) == 35, f"35 segnaposto nel template (trovati {len(ph)})")
check(ph == set(mapping), "i segnaposto del template coincidono col mapping")

print("== 4. Compilazione del template ==")
buf, stats = fill_pptx(BASE / "template_esempio.pptx", mapping)
check(stats["n_non_risolti"] == 0, f"0 segnaposto non risolti (={stats['n_non_risolti']})")
check(stats["n_sostituiti"] == 35, f"35 segnaposto sostituiti (={stats['n_sostituiti']})")
check(len(stats["inutilizzati"]) == 0, "0 voci di mapping inutilizzate")
filled = Presentation(buf)
txt = all_text(filled)
check("{{" not in txt and "}}" not in txt, "nessuna graffa residua nel file compilato")
check("2 - ROTTA" in txt and "1 - DERIVA" in txt, "risultati presenti nel testo")
check("inizia a definire" in txt, "commento strategia presente")
check("Integrare gli obiettivi DEIA" in txt, "attività sottogruppo presente")
# salva una copia compilata di esempio
buf.seek(0)
(BASE / "assessment_demo_compilato.pptx").write_bytes(buf.getvalue())

print("== 5. Robustezza: segnaposto spezzato su più run ==")
prs = Presentation()
sl = prs.slides.add_slide(prs.slide_layouts[6])
tb = sl.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
p = tb.text_frame.paragraphs[0]
for frag in ["{{strat", "egia.ris", "ultato}}"]:   # spezzato in 3 run
    r = p.add_run(); r.text = frag; r.font.size = Pt(18)
b2, s2 = fill_pptx(_save(prs), mapping)
t2 = all_text(Presentation(b2))
check(t2.strip() == "2 - ROTTA", f"run spezzati ricomposti -> '2 - ROTTA' (got '{t2.strip()}')")

print("== 6. Robustezza: segnaposto dentro una tabella ==")
prs = Presentation()
sl = prs.slides.add_slide(prs.slide_layouts[6])
gt = sl.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(6), Inches(1)).table
gt.cell(0, 0).text = "{{cultura.commento}}"
b3, s3 = fill_pptx(_save(prs), mapping)
t3 = all_text(Presentation(b3))
check("Valori, comportamenti" in t3 and "{{" not in t3, "segnaposto in tabella sostituito")

print("== 7. Radar dinamico sulla slide finale ==")
from pptx.enum.shapes import MSO_SHAPE_TYPE
buf.seek(0)
last = Presentation(buf)                       # template compilato al punto 4
final_slide = last.slides[len(last.slides._sldIdLst) - 1]
pics = [sh for sh in final_slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
check(len(pics) == 1, f"1 immagine radar sulla slide finale (trovate {len(pics)})")
radar_txt = all_text(last)
check("{{radar" not in radar_txt, "nessun marcatore radar residuo nel file compilato")

print("\nJSON valido:", end=" ")
json.loads((BASE / "deia_framework.json").read_text(encoding="utf-8"))
print("sì")

print("\nRISULTATO:", "TUTTO OK ✅" if ok else "CI SONO FALLIMENTI ❌")
raise SystemExit(0 if ok else 1)
