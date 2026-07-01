# -*- coding: utf-8 -*-
"""
Crea template_esempio.pptx : una slide per area, nello stile del documento,
con i SEGNAPOSTO {{...}} che l'app Streamlit andrà a compilare.

Serve per testare subito l'app (carichi questo + deia_mapping.xlsx).
Puoi sostituirlo con il TUO PowerPoint: basta usare gli stessi segnaposto.
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).resolve().parent
fw = json.loads((BASE / "deia_framework.json").read_text(encoding="utf-8"))

ORANGE = RGBColor(0xE8, 0x6A, 0x17)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_text(slide, l, t, w, h, runs, size, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, anchor=MSO_ANCHOR.TOP):
    """runs: stringa singola oppure lista di (testo, bold, color, size)."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, bold, color, size)]
    for txt, b, c, s in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Calibri"
    return box


# ----------------------------------------------------------------- SLIDE 1: cover
COVER = BASE / "cover.jpg"
if COVER.exists():
    cover = prs.slides.add_slide(blank)
    cover.shapes.add_picture(str(COVER), Inches(3.62), Inches(0.33),
                             Inches(5.43), Inches(6.85))


# ------------------------------------------- SLIDE 2: panoramica livello macro
def add_overview_slide():
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42), Inches(1.15), Inches(0.11))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    add_text(slide, 0.55, 0.52, 11.0, 0.8, "PANORAMICA", 30, bold=True)
    add_text(slide, 0.55, 1.18, 11.5, 0.4,
             "Livello di maturità DEIA aggregato sull'intero assessment", 13, color=GREY)

    # livello macro in evidenza
    add_text(slide, 0.55, 2.35, 6.0, 0.35, "LIVELLO MACRO", 13, bold=True)
    add_text(slide, 0.55, 2.78, 9.5, 1.5, "{{panoramica.risultato}}", 46, bold=True)
    add_text(slide, 0.55, 4.35, 6.0, 0.5, "{{panoramica.livello_medio}}", 18, color=GREY)

    # separatore
    ln = slide.shapes.add_shape(1, Inches(0.55), Inches(5.15), Inches(12.25), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()

    add_text(slide, 0.55, 5.35, 12.25, 0.35, "COSA SIGNIFICA", 13, bold=True)
    add_text(slide, 0.55, 5.75, 12.25, 1.5, "{{panoramica.descrizione}}", 12, color=GREY)


add_overview_slide()


for area in fw["aree"]:
    aid = area["id"]
    slide = prs.slides.add_slide(blank)

    # accento arancione + titolo area (statico)
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42), Inches(1.15), Inches(0.11))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    add_text(slide, 0.55, 0.52, 9.0, 0.8, area["nome"], 30, bold=True)

    # RISULTATO (label + segnaposto)
    add_text(slide, 0.55, 1.65, 2.3, 0.35, "RISULTATO", 13, bold=True)
    add_text(slide, 0.55, 2.02, 2.3, 0.4, "{{%s.risultato}}" % aid, 14, color=GREY)

    # COMMENTO (label + segnaposto)
    add_text(slide, 3.1, 1.65, 9.6, 0.35, "COMMENTO", 13, bold=True)
    add_text(slide, 3.1, 2.02, 9.7, 1.7, "{{%s.commento}}" % aid, 10.5,
             color=GREY)

    # intestazioni colonne sottogruppi
    y0 = 3.95
    add_text(slide, 0.55, y0, 3.4, 0.35, "SOTTOGRUPPO", 13, bold=True)
    add_text(slide, 4.15, y0, 2.2, 0.35, "RISULTATO", 13, bold=True)
    add_text(slide, 6.55, y0, 6.2, 0.35, "ATTIVITÀ SUGGERITE", 13, bold=True)

    # righe sottogruppi
    row_h = 1.02
    y = y0 + 0.45
    for sg in area["sottogruppi"]:
        sid = sg["id"]
        add_text(slide, 0.55, y, 3.45, row_h, sg["nome"], 12, bold=False)
        add_text(slide, 4.15, y, 2.2, row_h, "{{%s.risultato}}" % sid, 12, color=GREY)
        add_text(slide, 6.55, y, 6.25, row_h, "{{%s.attivita}}" % sid, 9.5, color=GREY)
        # sottile linea separatrice
        ln = slide.shapes.add_shape(1, Inches(0.55), Inches(y + row_h - 0.04),
                                    Inches(12.25), Emu(9525))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE
        ln.line.fill.background()
        y += row_h + 0.02


# ------------------------------------------------ SLIDE FINALE: profilo + radar
def add_profile_slide():
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42), Inches(1.15), Inches(0.11))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    add_text(slide, 0.55, 0.52, 11.0, 0.8, "PROFILO DI MATURITÀ", 30, bold=True)
    add_text(slide, 0.55, 1.18, 11.5, 0.4,
             "Distribuzione dei livelli di maturità per area", 13, color=GREY)

    # marcatore radar: alla compilazione diventa un'immagine dinamica
    add_text(slide, 0.7, 1.85, 6.7, 5.3, "{{radar.aree}}", 12, color=LINE)

    # colonna destra: livelli per area + livello macro
    x = 8.05
    add_text(slide, x, 1.95, 4.8, 0.35, "LIVELLI PER AREA", 13, bold=True)
    y = 2.45
    for area in fw["aree"]:
        add_text(slide, x, y, 4.8, 0.5,
                 [(area["nome"] + "   ", True, DARK, 13),
                  ("{{%s.risultato}}" % area["id"], False, GREY, 13)], 13)
        y += 0.58
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y + 0.05), Inches(4.65), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()
    add_text(slide, x, y + 0.22, 4.8, 0.35, "LIVELLO MACRO", 13, bold=True)
    add_text(slide, x, y + 0.62, 4.8, 0.5, "{{panoramica.risultato}}", 17, bold=True, color=ORANGE)


add_profile_slide()


out = BASE / "template_esempio.pptx"
prs.save(out)
print(f"Slide: {len(prs.slides._sldIdLst)} | Scritto: {out}")
